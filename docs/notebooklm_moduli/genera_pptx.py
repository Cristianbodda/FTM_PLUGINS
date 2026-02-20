"""
Genera 10 presentazioni PowerPoint per il video corso FTM Academy.
Ogni slide ha placeholder per screenshot con descrizione.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUT_DIR = os.path.dirname(__file__)
SCREENSHOTS_DIR = os.path.join(OUT_DIR, 'screenshots')

# Colori tema
PRIMARY = RGBColor(0x66, 0x7E, 0xEA)
SECONDARY = RGBColor(0x76, 0x4B, 0xA2)
DARK = RGBColor(0x2C, 0x3E, 0x50)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
GRAY = RGBColor(0x99, 0x99, 0x99)
GREEN = RGBColor(0x28, 0xA7, 0x45)
RED = RGBColor(0xDC, 0x35, 0x45)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)

def add_title_slide(prs, module_num, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PRIMARY
    # Title
    left, top, width, height = Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'Modulo {module_num}'
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xFF)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER
    # Subtitle
    left2, top2, width2, height2 = Inches(1.5), Inches(4), Inches(7), Inches(1)
    txBox2 = slide.shapes.add_textbox(left2, top2, width2, height2)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p3 = tf2.paragraphs[0]
    p3.text = subtitle
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)
    p3.alignment = PP_ALIGN.CENTER
    # Footer
    left3, top3 = Inches(0.8), Inches(6.5)
    txBox3 = slide.shapes.add_textbox(left3, top3, Inches(8.4), Inches(0.4))
    tf3 = txBox3.text_frame
    p4 = tf3.paragraphs[0]
    p4.text = 'FTM Academy - Video Corso Coach/Formatore'
    p4.font.size = Pt(10)
    p4.font.color.rgb = RGBColor(0xAA, 0xAA, 0xFF)
    p4.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets, note=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Title bar
    left, top, width, height = Inches(0), Inches(0), Inches(10), Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(15)
    # Left indent
    txBox.left = Inches(0.6)
    txBox.width = Inches(8.8)
    # Underline
    line_shape = slide.shapes.add_shape(
        1, Inches(0.6), Inches(0.95), Inches(8.8), Inches(0.04)  # MSO_SHAPE.RECTANGLE
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = PRIMARY
    line_shape.line.fill.background()
    # Bullets
    left2, top2, width2, height2 = Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.5)
    txBox2 = slide.shapes.add_textbox(left2, top2, width2, height2)
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f'  \u2022  {bullet}'
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
    # Note box
    if note:
        left3, top3 = Inches(0.8), Inches(6)
        txBox3 = slide.shapes.add_textbox(left3, top3, Inches(8.4), Inches(0.6))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p = tf3.paragraphs[0]
        p.text = f'\U0001f4a1 {note}'
        p.font.size = Pt(11)
        p.font.italic = True
        p.font.color.rgb = GRAY

def add_screenshot_slide(prs, title, screenshot_desc, screenshot_file=None):
    """Slide con placeholder per screenshot o screenshot reale."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(8.8), Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = DARK

    # Screenshot area
    img_left, img_top = Inches(0.5), Inches(0.9)
    img_width, img_height = Inches(9), Inches(5.5)

    screenshot_path = None
    if screenshot_file:
        full_path = os.path.join(SCREENSHOTS_DIR, screenshot_file)
        if os.path.exists(full_path):
            screenshot_path = full_path

    if screenshot_path:
        # Inserisci immagine reale
        slide.shapes.add_picture(screenshot_path, img_left, img_top, img_width)
    else:
        # Placeholder grigio
        placeholder = slide.shapes.add_shape(
            1, img_left, img_top, img_width, img_height
        )
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = LIGHT_GRAY
        placeholder.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        placeholder.line.width = Pt(1)
        # Testo placeholder
        tf2 = placeholder.text_frame
        tf2.word_wrap = True
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
        p1 = tf2.paragraphs[0]
        p1.text = '\U0001f4f7 SCREENSHOT'
        p1.font.size = Pt(24)
        p1.font.color.rgb = GRAY
        p1.font.bold = True
        p1.space_before = Pt(80)
        p2 = tf2.add_paragraph()
        p2.text = screenshot_desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf2.add_paragraph()
        p3.text = f'\n[File: {screenshot_file or "da generare"}]'
        p3.font.size = Pt(10)
        p3.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
        p3.alignment = PP_ALIGN.CENTER

    # Caption
    txBox3 = slide.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(8.8), Inches(0.5))
    tf3 = txBox3.text_frame
    p = tf3.paragraphs[0]
    p.text = screenshot_desc
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER

def add_end_slide(prs, module_num, next_title=''):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = SECONDARY
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f'Fine Modulo {module_num}'
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    if next_title:
        p2 = tf.add_paragraph()
        p2.text = f'\nProssimo: {next_title}'
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xFF)
        p2.alignment = PP_ALIGN.CENTER


# =============================================
# DEFINIZIONE 10 MODULI
# =============================================

moduli = [
    {
        'num': 1,
        'title': 'Introduzione e Accesso',
        'subtitle': 'Panoramica del sistema FTM Academy e come accedere',
        'next': 'Coach Dashboard - Panoramica, Viste e Filtri',
        'slides': [
            ('content', "Cos'\u00e8 FTM Academy", [
                'Ecosistema di plugin Moodle per competenze professionali',
                'Il coach monitora, valuta e guida gli studenti',
                'Percorso formativo strutturato in 6 settimane',
                'Dati oggettivi per colloqui e valutazioni',
            ]),
            ('screenshot', 'I 5 Strumenti del Coach', 'Schema dei 5 strumenti: Dashboard, Report, Valutazione, Bilancio, Scheduler', 'mod01_strumenti.png'),
            ('content', 'Le 4 Fonti Dati', [
                'Quiz: risultati oggettivi dei test Moodle (automatico)',
                'Autovalutazione: lo studente si valuta su scala Bloom (1-6)',
                'LabEval: valutazione pratica del formatore laboratorio',
                'Coach: la tua valutazione diretta basata su osservazione',
            ], 'Il radar overlay sovrappone tutte e 4 le fonti per confronto immediato'),
            ('screenshot', 'Accesso al Sistema', 'Pagina di login Moodle con URL e campi username/password', 'mod01_login.png'),
            ('screenshot', 'FTM Tools Hub', 'Blocco FTM Tools con i link ai vari strumenti coach', 'mod01_ftm_hub.png'),
            ('content', 'Flusso di Lavoro 6 Settimane', [
                'Sett. 1-2: Accoglienza, quiz iniziali, prima autovalutazione',
                'Sett. 3-4: Monitoraggio, laboratori, valutazioni intermedie',
                'Sett. 5-6: Colloqui, valutazione finale, report conclusivo',
                'Dashboard: centro di controllo per tutte le fasi',
            ]),
        ]
    },
    {
        'num': 2,
        'title': 'Coach Dashboard V2',
        'subtitle': 'Panoramica, Viste disponibili e Filtri avanzati',
        'next': 'Card Studente e Azioni Rapide',
        'slides': [
            ('screenshot', 'Dashboard V2 - Vista Completa', 'Vista panoramica della Coach Dashboard V2 con tutte le aree visibili', 'mod02_dashboard_full.png'),
            ('content', 'I 3 Pulsanti Header', [
                '\u2190 Versione Classica: torna al vecchio layout',
                'Scelte Rapide: assegnazione rapida test e laboratori',
                'Rapporto Classe: report aggregato di tutta la classe',
            ]),
            ('screenshot', 'Le 5 Statistiche', 'Le 5 stat card colorate: Studenti, Media, Autoval, Lab, Fine 6 Sett.', 'mod02_stats.png'),
            ('screenshot', 'I 6 Quick Filters', 'Pulsanti quick filter con conteggi: Tutti, Mancano Scelte, Manca Autoval...', 'mod02_quick_filters.png'),
            ('content', 'Le 4 Viste Disponibili', [
                'Classica: layout originale (retrocompatibilit\u00e0)',
                'Compatta: tabella una riga per studente (visione d\'insieme)',
                'Standard: card espandibili in griglia (consigliata)',
                'Dettagliata: pannelli sempre aperti, 2 colonne',
            ]),
            ('screenshot', 'Vista Compatta', 'Vista compatta con tabella a riga singola per studente', 'mod02_vista_compatta.png'),
            ('screenshot', 'Vista Standard', 'Vista standard con card espandibili, una aperta e due chiuse', 'mod02_vista_standard.png'),
            ('content', 'Zoom e Accessibilit\u00e0', [
                'A- (90%): testo pi\u00f9 piccolo, pi\u00f9 studenti visibili',
                'A (100%): dimensione normale',
                'A+ (120%): testo pi\u00f9 grande',
                'A++ (140%): massima accessibilit\u00e0',
            ], 'Le preferenze di zoom vengono salvate automaticamente'),
            ('screenshot', 'Filtri Avanzati', 'Pannello filtri con dropdown corso, chip colore gruppo, settimana e stato', 'mod02_filtri.png'),
        ]
    },
    {
        'num': 3,
        'title': 'Card Studente e Azioni',
        'subtitle': 'Tutto quello che trovi nella card di ogni studente',
        'next': 'Report Studente - Tab e Radar',
        'slides': [
            ('screenshot', 'Struttura Card Studente', 'Card studente espansa con tutte le sezioni visibili', 'mod03_card_completa.png'),
            ('content', 'Header della Card', [
                'Nome e email dello studente',
                'Badge settore con medaglie: \U0001f947 Primario, \U0001f948 Secondario, \U0001f949 Terziario',
                'Badge settimana corrente (es. "Sett. 3")',
                'Badge alert: "FINE 6 SETT." (rosso), "SOTTO SOGLIA" (arancione)',
            ]),
            ('content', 'Sezione Progress e Status', [
                '3 barre colorate: Competenze (viola), Autoval (teal), Lab (arancione)',
                'Rosso se sotto soglia 50%',
                'Badge stato: done (\u2713 verde), pending (\u23f1 giallo), missing (\u2717 rosso)',
                'Timeline 6 settimane con indicatori visivi',
            ]),
            ('screenshot', 'Timeline e Attivit\u00e0', 'Timeline 6 settimane con icone stato + tabella attivit\u00e0 settimanali', 'mod03_timeline.png'),
            ('content', 'Le 7 Quick Actions', [
                '\U0001f4ca Report \u2192 apre report studente dettagliato',
                '\u2705 Valutazione \u2192 apre valutazione formatore',
                '\U0001f464 Profilo \u2192 profilo utente Moodle',
                '\U0001f4ac Colloquio \u2192 bilancio competenze',
                '\U0001f4c4 Word \u2192 export documento (solo fine 6 sett.)',
                '\U0001f4e8 Sollecita \u2192 invia reminder autovalutazione',
                '\U0001f4be Salva \u2192 salva scelte settimanali',
            ]),
            ('screenshot', 'Quick Actions e Scelte', 'Barra pulsanti azione + dropdown scelte settimanali', 'mod03_actions.png'),
        ]
    },
    {
        'num': 4,
        'title': 'Report Studente',
        'subtitle': 'Tab, Radar e pannello FTM completo',
        'next': 'Gap Analysis e Spunti Colloquio',
        'slides': [
            ('screenshot', 'Header del Report', 'Header con foto studente, nome, email, corso e % globale su sfondo viola', 'mod04_header.png'),
            ('content', 'Le 7 Tab Principali', [
                '\U0001f4ca Panoramica: mappa competenze con radar e card aree',
                '\U0001f4da Piano: competenze in 4 categorie (Critico/Migliorare/Buono/Eccellente)',
                '\U0001f4dd Quiz: tabella quiz con tentativi e punteggi',
                '\U0001f4ca Autovalutazione \u2197: link esterno (nuova scheda)',
                '\U0001f52c Laboratorio \u2197: link esterno (nuova scheda)',
                '\U0001f4c8 Progressi: grafico andamento nel tempo',
                '\U0001f4cb Dettagli: tabella completa con valutazione inline',
            ]),
            ('screenshot', 'Tab Panoramica con Radar', 'Tab Panoramica con radar aree di competenza e riepilogo', 'mod04_panoramica.png'),
            ('screenshot', 'Tab Dettagli', 'Tabella competenze con filtri, ordinamento e badge Bloom cliccabili', 'mod04_dettagli.png'),
            ('content', 'Valutazione Inline', [
                'Clicca sul badge Bloom nella colonna "Valutazione"',
                'Dropdown con livelli N/O + 1-6',
                'Auto-save dopo 500ms (debounce)',
                'Toast conferma: "Salvando..." poi "\u2713 Salvato!"',
            ], 'Comodo per valutazioni rapide. Per valutazione completa usa la pagina dedicata.'),
            ('screenshot', 'Pannello FTM - 6 Tab', 'Pannello FTM con tab Settori, Ultimi 7gg, Configurazione, Progresso, Gap, Spunti', 'mod04_ftm_panel.png'),
            ('screenshot', 'Tab Ultimi 7 Giorni', 'Tabella quiz recenti con stato (Completato/In corso/Scaduto) e link Review', 'mod04_ultimi7gg.png'),
            ('screenshot', 'Configurazione e Soglie', 'Toggle visualizzazione + slider soglie Gap Analysis', 'mod04_config.png'),
            ('screenshot', 'Stampa Personalizzata', 'Modal stampa con checkbox sezioni, ordine, tono e filtro settore', 'mod04_stampa.png'),
        ]
    },
    {
        'num': 5,
        'title': 'Gap Analysis e Spunti',
        'subtitle': 'Analisi scostamenti e preparazione colloquio',
        'next': 'Valutazione Formatore e Scala Bloom',
        'slides': [
            ('content', 'Cos\'\u00e8 la Gap Analysis', [
                'Gap = Autovalutazione% - Performance Quiz%',
                'Sopravvalutazione: studente si valuta MEGLIO dei quiz (\u2b06\ufe0f rosso)',
                'Sottovalutazione: studente si valuta PEGGIO dei quiz (\u2b07\ufe0f arancione)',
                'Allineato: differenza sotto soglia (\u2705 verde)',
            ]),
            ('content', 'Soglie Configurabili', [
                'Allineamento: 5-40% (default 10%) \u2192 considerato allineato',
                'Monitorare: 15-60% (default 25%) \u2192 da tenere d\'occhio',
                'Critico: 20-80% (default 30%) \u2192 richiede intervento',
                'Ordinamento: per magnitudine gap (pi\u00f9 grande prima)',
            ]),
            ('screenshot', 'Tab Gap Analysis', 'Tabella Gap Analysis con indicatori freccia e colori per area', 'mod05_gap.png'),
            ('content', 'Tab Spunti Colloquio - 3 Categorie', [
                '\U0001f534 Critici: gap > soglia critico \u2192 domande specifiche urgenti',
                '\u26a0\ufe0f Attenzione: gap medio \u2192 suggerimenti miglioramento',
                '\u2705 Positivi: sottovalutazione \u2192 punti di forza da valorizzare',
            ]),
            ('screenshot', 'Spunti Colloquio', 'Le 3 categorie di spunti con domande suggerite per area', 'mod05_spunti.png'),
            ('content', 'Workflow Preparazione Colloquio', [
                '1. Apri Gap Analysis \u2192 identifica le aree critiche',
                '2. Leggi Spunti Colloquio \u2192 prepara le domande',
                '3. Scrivi note coach \u2192 annota punti da discutere',
                '4. (Opzionale) Stampa con tono formale per aziende',
            ]),
        ]
    },
    {
        'num': 6,
        'title': 'Valutazione e Scala Bloom',
        'subtitle': 'Come valutare con la scala Bloom (0-6)',
        'next': 'Bilancio Competenze e Note Coach',
        'slides': [
            ('screenshot', 'Pagina Valutazione Formatore', 'Header gradiente con selettore settore e medaglie', 'mod06_eval_header.png'),
            ('content', 'Struttura della Valutazione', [
                'Selezione settore: dropdown con \U0001f947\U0001f948\U0001f949 per priorit\u00e0',
                'Aree espandibili (A-G): accordion con competenze',
                '7 pulsanti per competenza: N/O, 1, 2, 3, 4, 5, 6',
                'Note per competenza + note generali',
            ]),
            ('screenshot', 'Compilazione Competenze', 'Area espansa con pulsanti Bloom e note per ogni competenza', 'mod06_eval_form.png'),
            ('content', 'Scala Bloom - I 6 Livelli', [
                '0 N/O: Non Osservato (non hai potuto valutare)',
                '1 Ricordare: sa elencare, definire, ripetere',
                '2 Comprendere: spiega con parole proprie il perch\u00e9',
                '3 Applicare: esegue procedure standard correttamente',
                '4 Analizzare: diagnostica, identifica cause',
                '5 Valutare: sceglie tra soluzioni e motiva',
                '6 Creare: inventa soluzioni nuove (raro in formazione)',
            ]),
            ('content', 'Esempi Pratici - Come Decidere', [
                'Ha ripetuto una definizione a memoria \u2192 Livello 1',
                'Ha spiegato un concetto con esempi propri \u2192 Livello 2',
                'Ha eseguito una procedura standard \u2192 Livello 3',
                'Ha identificato un guasto dai sintomi \u2192 Livello 4',
                'Ha scelto tra due soluzioni e motivato \u2192 Livello 5',
                'Ha inventato una soluzione inedita \u2192 Livello 6',
            ], 'Errore comune: confondere "Applicare" (segue la procedura da solo) con "Ricordare" (la ripete solo quando gliela detti)'),
            ('content', 'Salvataggio e Stati', [
                '\U0001f4be Salva Bozza (blu): puoi tornare a modificare',
                '\u2705 Salva e Completa (verde): valutazione finita',
                '\U0001f510 Firma (rosso): valutazione ufficiale e bloccata',
                '\U0001f513 Riapri per Modifiche (teal): sblocca una firmata',
            ]),
            ('screenshot', 'Pulsanti Azione', 'Barra azione con Salva Bozza, Completa, Firma e Riapri', 'mod06_eval_buttons.png'),
        ]
    },
    {
        'num': 7,
        'title': 'Bilancio Competenze e Note',
        'subtitle': 'Preparazione colloqui e sistema note',
        'next': 'Self-Assessment e FTM Scheduler',
        'slides': [
            ('content', 'Bilancio vs Report', [
                'Report Studente: analisi dati dettagliata (per il coach)',
                'Bilancio Competenze: sintesi per colloqui (per lo studente)',
                'Accesso: pulsante "\U0001f4ac Colloquio" dalla Dashboard',
                'URL: reports_v2.php?studentid=X',
            ]),
            ('screenshot', 'Bilancio - Tab Panoramica', '5 stat card cliccabili + punti di forza + aree critiche', 'mod07_bilancio_panoramica.png'),
            ('screenshot', 'Radar Confronto', 'Radar con checkbox Quiz/Autoval sovrapposti e grafico GAP', 'mod07_radar_confronto.png'),
            ('screenshot', 'Mappa Competenze', 'Card area cliccabili con colori stato (verde/teal/giallo/rosso)', 'mod07_mappa.png'),
            ('screenshot', 'Confronta Studenti', 'Due dropdown + tabella confronto + radar sovrapposto', 'mod07_confronta.png'),
            ('content', 'Tab Colloquio - Preparazione', [
                'Priorit\u00e0: aree ordinate per criticit\u00e0',
                'Domande suggerite per settore specifico',
                'Le 3 aree peggiori evidenziate',
                'Note coach con salvataggio AJAX',
            ]),
            ('content', 'Note Coach - Regole', [
                'Visibili a: te (coach) + segreteria',
                'MAI visibili allo studente',
                'Sovrascrittura: ogni salvataggio sostituisce il precedente',
                'Suggerimento: aggiungi la data prima di ogni annotazione',
            ]),
        ]
    },
    {
        'num': 8,
        'title': 'Self-Assessment e Scheduler',
        'subtitle': 'Gestione autovalutazioni e calendario formativo',
        'next': 'Casi d\'Uso Pratici',
        'slides': [
            ('screenshot', 'Self-Assessment Dashboard', 'Dashboard con 4 stat card + filtri + tabella studenti', 'mod08_selfassessment.png'),
            ('content', 'Azioni Self-Assessment', [
                '\U0001f441\ufe0f Vedi: apre il report autovalutazione',
                '\U0001f515 Disabilita/Riabilita: toggle AJAX con conferma',
                '\U0001f4e7 Reminder: invia promemoria (solo se abilitato + in attesa)',
                '\U0001f4cb Assegna Autovalutazioni: pagina di assegnazione',
            ]),
            ('screenshot', 'FTM Scheduler - Panoramica', 'Scheduler con barra gruppi attivi, stat card e tab', 'mod08_scheduler_full.png'),
            ('screenshot', 'Vista Settimanale', 'Griglia calendario 5 giorni x 2 fasce con blocchi colorati per gruppo', 'mod08_calendario.png'),
            ('content', 'I 5 Gruppi Colore', [
                '\U0001f7e8 Giallo (#FFFF00) - testo scuro',
                '\u2b1c Grigio (#808080) - testo bianco',
                '\U0001f7e5 Rosso (#FF0000) - testo bianco',
                '\U0001f7eb Marrone (#996633) - testo bianco',
                '\U0001f7ea Viola (#7030A0) - testo bianco',
            ], 'Bordo tratteggiato blu = prenotazioni esterne'),
            ('screenshot', 'Creare Nuovo Gruppo', 'Modal "Crea Nuovo Gruppo" con selettore colore e settimana', 'mod08_nuovo_gruppo.png'),
            ('screenshot', 'Tab Aule', 'Card delle 3 aule con tipo, capacit\u00e0 e attrezzature', 'mod08_aule.png'),
            ('screenshot', 'Tab Atelier', 'Catalogo atelier con tabella e riga obbligatoria evidenziata', 'mod08_atelier.png'),
            ('screenshot', 'Registro Presenze', 'Pagina presenze con card attivit\u00e0 e tabella studenti Presente/Assente', 'mod08_presenze.png'),
        ]
    },
    {
        'num': 9,
        'title': 'Casi d\'Uso Pratici',
        'subtitle': '16 workflow reali passo-passo',
        'next': 'Risoluzione Problemi e Riferimenti',
        'slides': [
            ('content', 'Caso 1: Nuovo Studente', [
                'Apri Dashboard \u2192 cerca Mario \u2192 espandi card',
                'Verifica settore e corso',
                'Assegna scelte settimanali (Test + Lab)',
                'Invia autovalutazione con pulsante Sollecita',
            ]),
            ('content', 'Caso 3: Gap Critici + Colloquio', [
                'Apri Report \u2192 tab FTM \u2192 Gap Analysis',
                'Identifica aree con gap > 30%',
                'Vai a Spunti Colloquio \u2192 leggi domande suggerite',
                'Apri Bilancio \u2192 tab Colloquio \u2192 salva note',
            ]),
            ('content', 'Caso 4: Valutazione Completa', [
                'Dashboard \u2192 Valutazione \u2192 seleziona settore',
                'Espandi area A \u2192 clicca livelli Bloom per ogni competenza',
                'Aggiungi note per competenza',
                'Salva bozza \u2192 Completa \u2192 Firma',
            ]),
            ('content', 'Caso 6: Fine Percorso + Export', [
                'Completa e firma la valutazione',
                'Dashboard \u2192 pulsante Word (appare solo a fine 6 sett.)',
                'Scarica il documento Word generato',
                'Verifica contenuto prima di inviare',
            ]),
            ('content', 'Caso 9: Colloquio Tecnico', [
                'Bilancio \u2192 tab Colloquio \u2192 le 3 aree peggiori',
                'Leggi domande suggerite specifiche per settore',
                'Scrivi note pre-colloquio',
                'Dopo il colloquio: aggiorna note con esito',
            ]),
            ('content', 'Caso 14: Filtri Combinati', [
                'Dashboard \u2192 Quick Filter "Sotto Soglia"',
                'Filtri Avanzati \u2192 Colore: Rosso + Settimana: 4',
                'Risultato: solo studenti critici del gruppo rosso sett. 4',
                'Per ognuno: Report \u2192 Gap Analysis \u2192 piano d\'azione',
            ]),
        ]
    },
    {
        'num': 10,
        'title': 'Problemi e Riferimenti',
        'subtitle': 'Troubleshooting, glossario e checklist',
        'next': '',
        'slides': [
            ('content', 'Problemi Comuni - Dashboard', [
                'Non vedo studenti: verifica corso e filtri attivi',
                'Pagina lenta: troppe card aperte \u2192 usa Vista Compatta',
                'Note perse: le note sovrascrivono, usa date nel testo',
                'Filtri non funzionano: clicca "Tutti" per reset completo',
            ]),
            ('content', 'Problemi Comuni - Report', [
                'Report vuoto: studente ha quiz? Verifica tab "Ultimi 7gg"',
                'Radar vuoti: servono almeno 3 competenze con dati',
                'Quiz non appare: verificare stato (completato vs in corso)',
                'Valutazione non salva: sessione scaduta \u2192 ricarica pagina',
            ]),
            ('content', 'Problemi Comuni - Altro', [
                'Export Word non funziona: solo per studenti "Fine 6 Sett."',
                'Sollecito non parte: studente gi\u00e0 completato o disabilitato',
                'Iscrizione Atelier fallisce: posti esauriti o non in sett. 3+',
                'Sessione scade: salva spesso, max ~2 ore di sessione',
            ]),
            ('content', 'Contatti Supporto', [
                'Studente non assegnato / permessi / corsi \u2192 Segreteria FTM',
                'Bug / errori pagina / funzionalit\u00e0 \u2192 Supporto tecnico IT',
                'Problemi Moodle generali (login, password) \u2192 Admin Moodle',
                'Domande sull\'uso degli strumenti \u2192 Coordinatore coach',
            ]),
            ('content', 'Checklist Settimanale', [
                'Luned\u00ec: nuovi studenti, fine 6 sett., sotto soglia, autoval mancanti',
                'Mar-Gio: aggiorna valutazioni, prendi note, verifica quiz, atelier',
                'Venerd\u00ec: rivedi sett. 5-6, firma valutazioni, export Word, report classe',
                'Mensile: confronta studenti, rivedi soglie, report per segreteria',
            ]),
            ('content', 'Riepilogo Corso Completo', [
                'Dashboard \u2192 il tuo centro di controllo quotidiano',
                'Report \u2192 analisi profonda quando serve approfondire',
                'Valutazione \u2192 la tua voce professionale (scala Bloom)',
                'Bilancio \u2192 preparazione colloqui basata su dati',
                'Scheduler \u2192 pianificazione e organizzazione',
            ], 'Buon lavoro! Ora hai tutti gli strumenti per essere un coach efficace.'),
        ]
    },
]

# =============================================
# GENERAZIONE
# =============================================

os.makedirs(SCREENSHOTS_DIR, exist_ok=True)

for mod in moduli:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    add_title_slide(prs, mod['num'], mod['title'], mod['subtitle'])

    # Content/Screenshot slides
    for slide_def in mod['slides']:
        if slide_def[0] == 'content':
            title = slide_def[1]
            bullets = slide_def[2]
            note = slide_def[3] if len(slide_def) > 3 else ''
            add_content_slide(prs, title, bullets, note)
        elif slide_def[0] == 'screenshot':
            title = slide_def[1]
            desc = slide_def[2]
            filename = slide_def[3] if len(slide_def) > 3 else None
            add_screenshot_slide(prs, title, desc, filename)

    # End slide
    add_end_slide(prs, mod['num'], mod['next'])

    # Save
    filename = f"PPT_Modulo_{mod['num']:02d}_{mod['title'].replace(' ', '_').replace('/', '_').replace(chr(39), '')}.pptx"
    filepath = os.path.join(OUT_DIR, filename)
    prs.save(filepath)
    slide_count = len(prs.slides)
    print(f"  OK {filename} ({slide_count} slide)")

print(f"\nTutti i 10 PowerPoint generati in: {OUT_DIR}")
